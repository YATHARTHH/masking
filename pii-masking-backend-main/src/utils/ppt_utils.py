import os
import asyncio
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
import tempfile
import shutil
from typing import List, Tuple
import platform
import subprocess

# Import the image processing function
from src.utils.image_utils import process_image

# Try to import Windows-specific libraries
IS_WINDOWS = platform.system() == 'Windows'
if IS_WINDOWS:
    try:
        import comtypes.client
        import win32com.client
        from pywintypes import com_error
        WINDOWS_COM_AVAILABLE = True
    except ImportError:
        WINDOWS_COM_AVAILABLE = False
        print("Warning: pywin32/comtypes not available. .ppt conversion will not work.")
else:
    WINDOWS_COM_AVAILABLE = False

UPLOAD_FOLDER = "uploads"
PROCESSED_FOLDER = "processed"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(PROCESSED_FOLDER, exist_ok=True)


def convert_ppt_to_pptx(ppt_path: str) -> str:
    """
    Convert .ppt file to .pptx format using available methods.
    
    Args:
        ppt_path: Path to the .ppt file
        
    Returns:
        Path to the converted .pptx file
    """
    if not IS_WINDOWS or not WINDOWS_COM_AVAILABLE:
        # Try LibreOffice as fallback
        return convert_ppt_to_pptx_libreoffice(ppt_path)
    
    import pythoncom
    pythoncom.CoInitialize()
    
    try:
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1
        
        # Open the .ppt file
        deck = powerpoint.Presentations.Open(os.path.abspath(ppt_path))
        
        # Create output path with .pptx extension
        pptx_path = ppt_path.rsplit('.', 1)[0] + '_converted.pptx'
        
        # Save as .pptx (format 24 is for .pptx)
        deck.SaveAs(os.path.abspath(pptx_path), 24)
        deck.Close()
        powerpoint.Quit()
        
        return pptx_path
    
    except Exception as e:
        print(f"Error converting .ppt to .pptx with COM: {e}")
        # Try LibreOffice as fallback
        return convert_ppt_to_pptx_libreoffice(ppt_path)
    finally:
        pythoncom.CoUninitialize()


def convert_ppt_to_pptx_libreoffice(ppt_path: str) -> str:
    """
    Convert .ppt to .pptx using LibreOffice (cross-platform).
    
    Args:
        ppt_path: Path to the .ppt file
        
    Returns:
        Path to the converted .pptx file
    """
    output_dir = os.path.dirname(ppt_path)
    
    try:
        # Try to use LibreOffice
        subprocess.run([
            'libreoffice',
            '--headless',
            '--convert-to', 'pptx',
            '--outdir', output_dir,
            ppt_path
        ], check=True, capture_output=True)
        
        pptx_path = ppt_path.rsplit('.', 1)[0] + '.pptx'
        return pptx_path
    
    except (subprocess.CalledProcessError, FileNotFoundError) as e:
        print(f"LibreOffice conversion failed: {e}")
        raise Exception("Cannot convert .ppt to .pptx. Please install LibreOffice or use Windows with MS PowerPoint.")


def pptx_to_images(pptx_path: str, output_dir: str) -> List[Tuple[str, Tuple[int, int]]]:
    """
    Convert PowerPoint slides to images using available methods.
    
    Args:
        pptx_path: Path to the PowerPoint file
        output_dir: Directory to save the images
        
    Returns:
        List of tuples containing (image_path, (width, height))
    """
    os.makedirs(output_dir, exist_ok=True)
    
    # Load presentation to get dimensions
    prs = Presentation(pptx_path)
    
    # Get slide dimensions in EMUs (English Metric Units)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # Convert EMUs to pixels (assuming 96 DPI)
    # 1 inch = 914400 EMUs, 1 inch = 96 pixels at 96 DPI
    width_px = int(slide_width / 914400 * 96)
    height_px = int(slide_height / 914400 * 96)
    
    # Try Windows COM method first
    if IS_WINDOWS and WINDOWS_COM_AVAILABLE:
        try:
            return pptx_to_images_com(pptx_path, output_dir, width_px, height_px)
        except Exception as e:
            print(f"COM export failed, trying LibreOffice: {e}")
    
    # Fallback to LibreOffice
    return pptx_to_images_libreoffice(pptx_path, output_dir, width_px, height_px)


def pptx_to_images_com(pptx_path: str, output_dir: str, width_px: int, height_px: int) -> List[Tuple[str, Tuple[int, int]]]:
    """
    Export slides using COM automation (Windows only).
    """
    import pythoncom
    pythoncom.CoInitialize()
    
    try:
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1
        
        deck = powerpoint.Presentations.Open(os.path.abspath(pptx_path))
        
        image_paths = []
        for i, slide in enumerate(deck.Slides):
            image_path = os.path.join(output_dir, f"slide_{i+1}.png")
            
            # Export slide as image (PNG format = 17)
            slide.Export(os.path.abspath(image_path), "PNG", width_px, height_px)
            
            image_paths.append((image_path, (width_px, height_px)))
        
        deck.Close()
        powerpoint.Quit()
        
        return image_paths
        
    finally:
        pythoncom.CoUninitialize()


def pptx_to_images_libreoffice(pptx_path: str, output_dir: str, width_px: int, height_px: int) -> List[Tuple[str, Tuple[int, int]]]:
    """
    Export slides using LibreOffice (cross-platform).
    """
    try:
        # Convert to PDF first
        pdf_path = os.path.join(output_dir, "temp.pdf")
        subprocess.run([
            'libreoffice',
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', output_dir,
            pptx_path
        ], check=True, capture_output=True)
        
        # Rename the PDF
        expected_pdf = os.path.join(output_dir, os.path.basename(pptx_path).rsplit('.', 1)[0] + '.pdf')
        if os.path.exists(expected_pdf):
            os.rename(expected_pdf, pdf_path)
        
        # Convert PDF to images using pdf2image
        from pdf2image import convert_from_path
        
        images = convert_from_path(pdf_path, dpi=150)
        
        image_paths = []
        for i, img in enumerate(images):
            image_path = os.path.join(output_dir, f"slide_{i+1}.png")
            
            # Resize to match expected dimensions
            img_resized = img.resize((width_px, height_px), Image.LANCZOS)
            img_resized.save(image_path, 'PNG')
            
            image_paths.append((image_path, (width_px, height_px)))
        
        # Clean up PDF
        if os.path.exists(pdf_path):
            os.remove(pdf_path)
        
        return image_paths
        
    except (subprocess.CalledProcessError, FileNotFoundError) as e:
        print(f"LibreOffice conversion failed: {e}")
        raise Exception("Cannot convert slides to images. Please install LibreOffice.")


def images_to_pptx(image_paths: List[str], output_path: str, slide_dimensions: Tuple[int, int]):
    """
    Create a PowerPoint presentation from images.
    
    Args:
        image_paths: List of paths to masked images
        output_path: Path where the output PowerPoint file will be saved
        slide_dimensions: Tuple of (width, height) in pixels
    """
    # Create a new presentation
    prs = Presentation()
    
    # Set slide dimensions (convert pixels back to EMUs)
    width_px, height_px = slide_dimensions
    prs.slide_width = int(width_px * 914400 / 96)
    prs.slide_height = int(height_px * 914400 / 96)
    
    # Use blank slide layout
    blank_slide_layout = prs.slide_layouts[6]  # 6 is typically the blank layout
    
    for image_path in image_paths:
        # Add a blank slide
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add the image to fill the entire slide
        left = top = Inches(0)
        
        # Calculate dimensions in inches
        width_inches = width_px / 96
        height_inches = height_px / 96
        
        slide.shapes.add_picture(
            image_path,
            left,
            top,
            width=Inches(width_inches),
            height=Inches(height_inches)
        )
    
    # Save the presentation
    prs.save(output_path)


async def process_ppt(ppt_path: str, pii_category: str, highlight_mode: str) -> str:
    """
    Process PowerPoint file by converting to images, masking PII, and converting back.
    
    Args:
        ppt_path: Path to the PowerPoint file (.ppt or .pptx)
        pii_category: Categories of PII to detect
        highlight_mode: Mode of masking ('blurring' or 'rectangular_box')
        
    Returns:
        Path to the processed PowerPoint file
    """
    # Create temporary directory for intermediate files
    temp_dir = tempfile.mkdtemp()
    
    try:
        # Get file extension
        file_ext = os.path.splitext(ppt_path)[1].lower()
        
        # Convert .ppt to .pptx if needed
        if file_ext == '.ppt':
            print("Converting .ppt to .pptx...")
            pptx_path = await asyncio.to_thread(convert_ppt_to_pptx, ppt_path)
        else:
            pptx_path = ppt_path
        
        # Step 1: Convert PowerPoint slides to images
        print("Converting slides to images...")
        slides_dir = os.path.join(temp_dir, "slides")
        image_data = await asyncio.to_thread(pptx_to_images, pptx_path, slides_dir)
        
        if not image_data:
            raise Exception("No slides found in PowerPoint file")
        
        # Get slide dimensions from the first slide
        _, slide_dimensions = image_data[0]
        
        # Step 2: Process each image to mask PII
        print("Masking PII in slides...")
        masked_images = []
        
        for image_path, _ in image_data:
            # Process the image using the image_utils function
            masked_path = await asyncio.to_thread(
                process_image,
                image_path,
                pii_category,
                highlight_mode,
                facial=False
            )
            masked_images.append(masked_path)
        
        # Step 3: Convert masked images back to PowerPoint
        print("Creating masked PowerPoint file...")
        output_filename = f"masked_{os.path.basename(ppt_path).rsplit('.', 1)[0]}.pptx"
        output_path = os.path.join(PROCESSED_FOLDER, output_filename)
        
        await asyncio.to_thread(images_to_pptx, masked_images, output_path, slide_dimensions)
        
        print(f"Successfully processed PowerPoint file: {output_path}")
        return output_path
    
    except Exception as e:
        print(f"Error processing PowerPoint: {e}")
        raise
    
    finally:
        # Clean up temporary directory
        try:
            shutil.rmtree(temp_dir)
            
            # Clean up converted .pptx file if it was created
            if file_ext == '.ppt' and os.path.exists(pptx_path):
                os.remove(pptx_path)
        except Exception as e:
            print(f"Error cleaning up temporary files: {e}")